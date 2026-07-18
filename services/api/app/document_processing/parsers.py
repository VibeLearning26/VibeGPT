"""
VibeGPT – Document Parsers

Parse PDF, PPTX, DOCX, XLSX files into structured content
with page/slide/sheet metadata for chunking.
"""

from __future__ import annotations

import io
import re
from dataclasses import dataclass
from typing import List, Optional

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches


@dataclass
class ParsedElement:
    """A single extractable element from a document."""
    content: str
    page_number: Optional[int] = None
    slide_number: Optional[int] = None
    sheet_name: Optional[str] = None
    heading: Optional[str] = None
    metadata: dict = None

    def __post_init__(self):
        if self.metadata is None:
            self.metadata = {}


def parse_pdf(file_bytes: bytes) -> List[ParsedElement]:
    """Parse PDF using PyMuPDF, preserving page numbers and headings."""
    elements = []
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    
    for page_num, page in enumerate(doc, start=1):
        # Extract text with structure
        text = page.get_text("text")
        if not text.strip():
            continue
        
        # Try to detect headings (first line, larger font, etc.)
        blocks = page.get_text("dict")["blocks"]
        heading = None
        for block in blocks:
            if "lines" in block:
                for line in block["lines"]:
                    for span in line["spans"]:
                        if span["size"] > 14 and span["text"].strip():
                            heading = span["text"].strip()
                            break
                    if heading:
                        break
            if heading:
                break
        
        # Clean up text
        text = clean_text(text)
        if text:
            elements.append(ParsedElement(
                content=text,
                page_number=page_num,
                heading=heading,
                metadata={"source": "pdf"}
            ))
    
    doc.close()
    return elements


def parse_pptx(file_bytes: bytes) -> List[ParsedElement]:
    """Parse PPTX using python-pptx, preserving slide numbers and titles."""
    elements = []
    prs = Presentation(io.BytesIO(file_bytes))
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = []
        heading = None
        
        # Check for title shape
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and not heading:
                    heading = text
                elif text:
                    slide_text.append(text)
        
        full_text = "\n".join(slide_text)
        full_text = clean_text(full_text)
        
        if full_text or heading:
            elements.append(ParsedElement(
                content=full_text,
                slide_number=slide_num,
                heading=heading,
                metadata={"source": "pptx"}
            ))
    
    return elements


def parse_docx(file_bytes: bytes) -> List[ParsedElement]:
    """Parse DOCX using python-docx, preserving headings and paragraphs."""
    elements = []
    doc = DocxDocument(io.BytesIO(file_bytes))
    
    current_heading = None
    current_content = []
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        
        # Detect heading styles
        if para.style.name.startswith("Heading"):
            # Save previous section
            if current_content:
                elements.append(ParsedElement(
                    content=clean_text("\n".join(current_content)),
                    heading=current_heading,
                    metadata={"source": "docx"}
                ))
                current_content = []
            current_heading = text
        else:
            current_content.append(text)
    
    # Save last section
    if current_content:
        elements.append(ParsedElement(
            content=clean_text("\n".join(current_content)),
            heading=current_heading,
            metadata={"source": "docx"}
        ))
    
    return elements


def parse_xlsx(file_bytes: bytes) -> List[ParsedElement]:
    """Parse XLSX using openpyxl, preserving sheet names and row structure."""
    elements = []
    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_text = []
        heading = None
        
        for row_idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
            row_values = [str(cell) for cell in row if cell is not None]
            if not row_values:
                continue
            
            # First non-empty row as heading
            if heading is None:
                heading = " | ".join(row_values)
            else:
                sheet_text.append(" | ".join(row_values))
        
        if sheet_text or heading:
            elements.append(ParsedElement(
                content=clean_text("\n".join(sheet_text)),
                sheet_name=sheet_name,
                heading=heading,
                metadata={"source": "xlsx"}
            ))
    
    wb.close()
    return elements


def clean_text(text: str) -> str:
    """Clean extracted text: normalize whitespace, remove artifacts."""
    if not text:
        return ""
    # Replace multiple newlines with double newline
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Replace multiple spaces with single space
    text = re.sub(r" {2,}", " ", text)
    # Remove zero-width and control chars
    text = re.sub(r"[\u200b-\u200f\ufeff]", "", text)
    return text.strip()


def parse_document(file_bytes: bytes, mime_type: str | SourceType) -> List[ParsedElement]:
    """Dispatch to appropriate parser based on MIME type or SourceType."""
    from app.models.document import SourceType
    
    # Convert SourceType enum to MIME type if needed
    if isinstance(mime_type, SourceType):
        source_to_mime = {
            SourceType.PDF_NOTES: "application/pdf",
            SourceType.PPTX_PRESENTATION: "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            SourceType.DOCX_NOTES: "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            SourceType.XLSX_QUESTION_BANK: "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        }
        mime_type = source_to_mime.get(mime_type)
        if not mime_type:
            raise ValueError(f"Unsupported source type: {mime_type}")
    
    mime_map = {
        "application/pdf": parse_pdf,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": parse_pptx,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": parse_docx,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": parse_xlsx,
    }
    
    parser = mime_map.get(mime_type)
    if not parser:
        raise ValueError(f"Unsupported MIME type: {mime_type}")
    
    return parser(file_bytes)