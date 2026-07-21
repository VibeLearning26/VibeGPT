"""
VibeGPT – Document Text Extraction

Extracts text units (with page/slide/sheet locations) from uploaded files.
Supported: PDF (PyMuPDF), PPTX (python-pptx), DOCX (python-docx), XLSX (openpyxl).
"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from pathlib import Path

logger = logging.getLogger(__name__)


class ExtractionError(Exception):
    """Raised when text cannot be extracted from a document."""


@dataclass
class TextUnit:
    """A block of extracted text with its location in the source document."""

    text: str
    page_number: int | None = None
    slide_number: int | None = None
    sheet_name: str | None = None
    heading: str | None = None


def extract_pdf(path: Path) -> list[TextUnit]:
    import fitz  # PyMuPDF

    units: list[TextUnit] = []
    try:
        with fitz.open(path) as doc:
            for i, page in enumerate(doc, start=1):
                text = page.get_text("text").strip()
                if text:
                    units.append(TextUnit(text=text, page_number=i))
    except Exception as e:
        raise ExtractionError(f"Failed to extract PDF text: {e}") from e
    return units


def extract_pptx(path: Path) -> list[TextUnit]:
    from pptx import Presentation

    units: list[TextUnit] = []
    try:
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            title = None
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if shape == getattr(slide.shapes, "title", None):
                    title = text
                texts.append(text)
            if texts:
                units.append(
                    TextUnit(text="\n".join(texts), slide_number=i, heading=title)
                )
    except Exception as e:
        raise ExtractionError(f"Failed to extract PPTX text: {e}") from e
    return units


def extract_docx(path: Path) -> list[TextUnit]:
    import docx

    units: list[TextUnit] = []
    try:
        document = docx.Document(str(path))
        current_heading: str | None = None
        buffer: list[str] = []

        def flush():
            if buffer:
                units.append(
                    TextUnit(text="\n".join(buffer), heading=current_heading)
                )
                buffer.clear()

        for para in document.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            if para.style.name.startswith("Heading"):
                flush()
                current_heading = text
            buffer.append(text)
        flush()
    except Exception as e:
        raise ExtractionError(f"Failed to extract DOCX text: {e}") from e
    return units


def extract_xlsx(path: Path) -> list[TextUnit]:
    import openpyxl

    units: list[TextUnit] = []
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                units.append(TextUnit(text="\n".join(rows), sheet_name=sheet.title))
        wb.close()
    except Exception as e:
        raise ExtractionError(f"Failed to extract XLSX text: {e}") from e
    return units


_EXTRACTORS = {
    ".pdf": extract_pdf,
    ".pptx": extract_pptx,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
}


def extract_text(path: Path) -> list[TextUnit]:
    """Extract text units from a file based on its extension."""
    extractor = _EXTRACTORS.get(path.suffix.lower())
    if extractor is None:
        raise ExtractionError(f"Unsupported file type: {path.suffix}")
    units = extractor(path)
    if not units:
        raise ExtractionError("No text could be extracted from the document")
    return units
